#!/usr/bin/env python3
"""
Generate a professional 16:9 PowerPoint Presentation (.pptx) with 19 slides based on the detailed report sections:
- Background, Problem Statement, Objectives, Scope, Target Audience
- Literature review, NIST password rules, Session flags, CSRF, Rate limits
- Normalized 3NF Database Design, SQLite/MySQL adaptability, parameterized queries
- Columns table schema for students & admins, SQL queries/CRUD write & read statements
- Frontend Apple HIG design system variables, 4-step wizard logic, local asset self-hosting CSP bypass
- Student & Admin output screens (Canvas particles, SVG rings, Chart.js KPIs, CSV downloads)
- Advantages/Limitations, final conclusions, future roadmap stages, and key references.
- Sapthagiri University logo on the top right of all pages matching the white background.
- Names and SRNs of all 3 members on all slides.
"""

import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Brand Colors
NAVY = RGBColor(0, 49, 126)      # #00317e
GOLD = RGBColor(212, 175, 55)    # #D4AF37
DARK_GRAY = RGBColor(28, 28, 30) # #1C1C1E
LIGHT_GRAY = RGBColor(242, 242, 247) # #F2F2F7
WHITE = RGBColor(255, 255, 255)  # #FFFFFF
MUTED_GRAY = RGBColor(110, 110, 115) # #6E6E73
SUCCESS_GREEN = RGBColor(16, 185, 129) # #10B981
WARNING_RED = RGBColor(239, 68, 68) # #EF4444

LOGO_PATH = "static/logo.png"

def create_presentation():
    # Initialize Presentation
    prs = Presentation()
    
    # Set to widescreen 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6] # Blank slide
    
    # 👥 Team Members info
    members_text = "Keerthana GP (24SUUBECS0914)  |  Keerthana BS (23SUUBECS0909)  |  Kiran M Biradar (24SUUBECS0937)"
    
    # Helper: Set slide background to white (so logo matches background)
    def set_white_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE

    # Helper: Add Logo at top right
    def add_logo(slide, is_title_slide=False):
        if os.path.exists(LOGO_PATH):
            if is_title_slide:
                # Place larger logo on title slide
                slide.shapes.add_picture(LOGO_PATH, Inches(4.5), Inches(1.5), width=Inches(4.33))
            else:
                # Standard top right logo
                slide.shapes.add_picture(LOGO_PATH, Inches(10.5), Inches(0.25), width=Inches(2.5))

    # Helper: Add Footer on content slides
    def add_footer(slide):
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.4))
        tf = footer_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.text = f"Student Hub  |  Sapthagiri University  |  {members_text}"
        p.font.name = "Arial"
        p.font.size = Pt(9.5)
        p.font.color.rgb = MUTED_GRAY
        p.alignment = PP_ALIGN.CENTER

    # Helper: Add standard slide template
    def add_content_slide(title_text):
        slide = prs.slides.add_slide(blank_layout)
        set_white_background(slide)
        add_logo(slide)
        add_footer(slide)
        
        # Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(9.5), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Arial"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = NAVY
        
        # Gold Accent Line below title
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(11.7), Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = GOLD
        line.line.color.rgb = GOLD
        
        return slide

    # =========================================================================
    # SLIDE 1: Title Slide (Introduction 1/2)
    # =========================================================================
    slide_1 = prs.slides.add_slide(blank_layout)
    set_white_background(slide_1)
    add_logo(slide_1, is_title_slide=True)
    
    # Title & Subtitle Card (Bottom half of title slide)
    title_box = slide_1.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11.7), Inches(1.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
    
    p1 = tf.paragraphs[0]
    p1.text = "STUDENT HUB"
    p1.font.name = "Arial"
    p1.font.size = Pt(48)
    p1.font.bold = True
    p1.font.color.rgb = NAVY
    p1.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "Online Student Registration System"
    p2.font.name = "Arial"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = GOLD
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)

    # Department and CSE label
    dept_box = slide_1.shapes.add_textbox(Inches(0.8), Inches(4.7), Inches(11.7), Inches(0.4))
    tf_dept = dept_box.text_frame
    tf_dept.word_wrap = True
    p_dept = tf_dept.paragraphs[0]
    p_dept.text = "Department of Computer Science & Engineering (Dept. of CSE)"
    p_dept.font.name = "Arial"
    p_dept.font.size = Pt(12)
    p_dept.font.italic = True
    p_dept.font.color.rgb = NAVY
    p_dept.alignment = PP_ALIGN.CENTER

    # Presenter / Author Details
    presenter_box = slide_1.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.8))
    tf_pres = presenter_box.text_frame
    tf_pres.word_wrap = True
    tf_pres.margin_left = tf_pres.margin_top = tf_pres.margin_bottom = tf_pres.margin_right = 0
    
    p_pres_title = tf_pres.paragraphs[0]
    p_pres_title.text = "DEVELOPMENT TEAM"
    p_pres_title.font.name = "Arial"
    p_pres_title.font.size = Pt(11)
    p_pres_title.font.bold = True
    p_pres_title.font.color.rgb = MUTED_GRAY
    p_pres_title.alignment = PP_ALIGN.CENTER
    
    team = [
        ("Keerthana GP", "24SUUBECS0914"),
        ("Keerthana BS", "23SUUBECS0909"),
        ("Kiran M Biradar", "24SUUBECS0937")
    ]
    
    for name, srn in team:
        p_member = tf_pres.add_paragraph()
        p_member.text = f"{name}  ({srn})"
        p_member.font.name = "Arial"
        p_member.font.size = Pt(13)
        p_member.font.bold = True
        p_member.font.color.rgb = DARK_GRAY
        p_member.alignment = PP_ALIGN.CENTER
        p_member.space_before = Pt(4)

    # =========================================================================
    # SLIDE 2: Background (Introduction 2/2)
    # =========================================================================
    slide_2 = add_content_slide("Introduction: Background & Concept")
    
    left_box = slide_2.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(7.5), Inches(4.8))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    tf_left.margin_left = tf_left.margin_top = tf_left.margin_bottom = tf_left.margin_right = 0
    
    p = tf_left.paragraphs[0]
    p.text = "Traditional student registration is traditionally a paper-intensive, time-consuming procedure requiring physical submission, manual entry, and long processing cycles."
    p.font.name = "Arial"
    p.font.size = Pt(17)
    p.font.color.rgb = DARK_GRAY
    p.space_after = Pt(14)
    
    bg_pts = [
        "📄 Paper Limitations: Transcription errors (1-3% per field), storage bulk, and long 3-5 days wait periods.",
        "🖥️ Legacy Systems Issues: Early electronic systems served as basic forms lacking real-time validation or mobile viewport support.",
        "⚡ Digitalization Pressure: Modern institutions face increasing pressure to digitize operations and meet student expectations.",
        "🛡️ Core Solution: Student Hub balances accessibility with multi-layered secure endpoints, real-time client indicators, and immediate ID generation."
    ]
    for pt in bg_pts:
        p = tf_left.add_paragraph()
        p.text = pt
        p.font.name = "Arial"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(10)

    # Right Card: Target Objectives
    right_card = slide_2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.4))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = LIGHT_GRAY
    right_card.line.color.rgb = GOLD
    right_card.line.width = Pt(1.5)
    
    rc_text = slide_2.shapes.add_textbox(Inches(9.1), Inches(2.1), Inches(3.1), Inches(3.8))
    tf_rc = rc_text.text_frame
    tf_rc.word_wrap = True
    p = tf_rc.paragraphs[0]
    p.text = "TARGET AUDIENCES"
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.space_after = Pt(14)
    
    audiences = [
        ("Students", "Expect fast (under 3 mins) enrollment and instant ID generation."),
        ("Enrollment Staff", "Need efficient tools. Reduces registration workloads by ~70%."),
        ("IT Administrators", "Require secure configurations (CSP, HSTS, Docker)."),
        ("Institutional Leaders", "Requires real-time dashboard visibility and audit trails.")
    ]
    for title, desc in audiences:
        p = tf_rc.add_paragraph()
        p.text = f"• {title}"
        p.font.name = "Arial"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = GOLD
        p.space_before = Pt(6)
        
        p = tf_rc.add_paragraph()
        p.text = desc
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY

    # =========================================================================
    # SLIDE 3: Problem Definition
    # =========================================================================
    slide_3 = add_content_slide("Problem Definition")
    
    problems = [
        (" P1: Excessive Completion Time", "Registration takes 30-45 minutes. Single-page scrolling forms with 20+ fields cause high cognitive fatigue.", Inches(0.8), Inches(1.8), "FO1: Limit application time to < 3 mins."),
        (" P2: Vulnerable Server Architectures", "Legacy databases lack encryption and rate limits. Plaintext password storage and missing CSRF tokens risk breaches.", Inches(6.8), Inches(1.8), "SO1-4: Implement Bcrypt (work factor 12) & CSRF check."),
        (" P3: Desktop-Only Interfaces", "Desktop viewports required. Small touch targets (<44px) make mobile enrollment extremely frustrating.", Inches(0.8), Inches(4.3), "UXO1-2: Mobile-first responsive views, 44px+ targets."),
        (" P4: Regulatory Compliance Gaps", "System lack log monitors (FERPA, GDPR) to track user updates, creating auditing compliance risks.", Inches(6.8), Inches(4.3), "AO1-4: Audit trails logger and dynamic charts dashboard.")
    ]
    
    for title, desc, left, top, objective in problems:
        card = slide_3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.7), Inches(2.1))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = WARNING_RED
        card.line.width = Pt(1)
        
        tb = slide_3.shapes.add_textbox(left + Inches(0.2), top + Inches(0.1), Inches(5.3), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Arial"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.space_after = Pt(4)
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.name = "Arial"
        p.font.size = Pt(12.5)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)
        
        p = tf.add_paragraph()
        p.text = f"🎯 Objective: {objective}"
        p.font.name = "Arial"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = SUCCESS_GREEN

    # =========================================================================
    # SLIDE 4: Database Design: ER Model (Database Design 1/2)
    # =========================================================================
    slide_4 = add_content_slide("Database Design: ER Model Concepts")
    
    left_box = slide_4.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(7.5), Inches(4.8))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    tf_left.margin_left = tf_left.margin_top = tf_left.margin_bottom = tf_left.margin_right = 0
    
    p = tf_left.paragraphs[0]
    p.text = "The transactional system models candidate profiles and system operations cleanly using relational Entity-Relationship standards."
    p.font.name = "Arial"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GRAY
    p.space_after = Pt(20)
    
    db_concepts = [
        "👤 Students Entity: Captures full applicant metrics (personal details, mailing addresses, school history, hashes).",
        "🛡️ Admins Entity: Stores administrator account records separated from standard student scopes.",
        "📊 Uniqueness Constraints: Enforces unique parameters on `email` and `student_id` database columns.",
        "📑 Indexing Configurations: Creates indexes (`idx_email`, `idx_student_id`) to optimize login search queries."
    ]
    for pt in db_concepts:
        p = tf_left.add_paragraph()
        p.text = pt
        p.font.name = "Arial"
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(12)

    # Right Card: Normalization
    right_card = slide_4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.4))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = LIGHT_GRAY
    right_card.line.color.rgb = GOLD
    right_card.line.width = Pt(1.5)
    
    rc_text = slide_4.shapes.add_textbox(Inches(9.1), Inches(2.1), Inches(3.1), Inches(3.8))
    tf_rc = rc_text.text_frame
    tf_rc.word_wrap = True
    p = tf_rc.paragraphs[0]
    p.text = "NORMALIZATION (3NF)"
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.space_after = Pt(16)
    
    normal_rules = [
        ("First Normal Form (1NF)", "Eliminates repeating columns and multi-valued attributes."),
        ("Second Normal Form (2NF)", "Removes partial dependencies on composite primary keys."),
        ("Third Normal Form (3NF)", "Removes transitive dependencies, securing data integrity.")
    ]
    for m_title, m_desc in normal_rules:
        p = tf_rc.add_paragraph()
        p.text = f"✔ {m_title}"
        p.font.name = "Arial"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = GOLD
        p.space_before = Pt(8)
        
        p = tf_rc.add_paragraph()
        p.text = m_desc
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY

    # =========================================================================
    # SLIDE 5: Database Architecture (Database Design 2/2)
    # =========================================================================
    slide_5 = add_content_slide("Database Architecture & SQL Injection Control")
    
    left_box = slide_5.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(6.5), Inches(4.8))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    tf_left.margin_left = tf_left.margin_top = tf_left.margin_bottom = tf_left.margin_right = 0
    
    p = tf_left.paragraphs[0]
    p.text = "The application implements a dual-database connectivity strategy, abstracting SQL query generation entirely through SQLAlchemy."
    p.font.name = "Arial"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GRAY
    p.space_after = Pt(14)
    
    arch_pts = [
        "🗄️ Dev Fallback (SQLite): Automatically creates a local database file `student_hub.db` if no MySQL settings exist.",
        "🌐 Prod Database (MySQL): Easily connects to MySQL 8+ schemas. URL-encodes special characters in passwords to prevent connection crashes.",
        "🛡️ SQL Injection Prevention: Uses SQLAlchemy ORM parameterized queries to separate user input variables from SQL commands, blocking injection attacks.",
        "🔑 Collision-Resistant IDs: Generates IDs in `STU-{YEAR}-{RANDOM4}` format, dynamically checking database availability via a collision-retry loop."
    ]
    for pt in arch_pts:
        p = tf_left.add_paragraph()
        p.text = pt
        p.font.name = "Arial"
        p.font.size = Pt(14.5)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(12)

    # Right Box: Model Sample
    right_card = slide_5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(1.8), Inches(4.7), Inches(4.4))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = DARK_GRAY
    right_card.line.color.rgb = GOLD
    right_card.line.width = Pt(1)
    
    rc_text = slide_5.shapes.add_textbox(Inches(7.95), Inches(2.0), Inches(4.4), Inches(4.0))
    tf_rc = rc_text.text_frame
    tf_rc.word_wrap = True
    p = tf_rc.paragraphs[0]
    p.text = "SQLALCHEMY ORM SCHEMAS"
    p.font.name = "Arial"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.space_after = Pt(12)
    
    model_code = (
        "class Student(db.Model, UserMixin):\n"
        "    __tablename__ = 'students'\n"
        "    id = db.Column(db.Integer, primary_key=True)\n"
        "    student_id = db.Column(\n"
        "        db.String(20), unique=True, nullable=False\n"
        "    )\n"
        "    email = db.Column(\n"
        "        db.String(120), unique=True, nullable=False\n"
        "    )\n"
        "    major = db.Column(\n"
        "        db.String(100), nullable=False\n"
        "    )\n"
        "    password_hash = db.Column(\n"
        "        db.String(255), nullable=False\n"
        "    )\n"
        "    registration_date = db.Column(\n"
        "        db.DateTime, default=datetime.utcnow\n"
        "    )"
    )
    p = tf_rc.add_paragraph()
    p.text = model_code
    p.font.name = "Courier New"
    p.font.size = Pt(9.5)
    p.font.color.rgb = SUCCESS_GREEN
    p.line_spacing = 1.15

    # =========================================================================
    # SLIDE 6: Tables & Relationships: Students Table Schema (Tables/Relationships 1/2)
    # =========================================================================
    slide_6 = add_content_slide("Tables & Relationships: Students Schema")
    
    # Table of columns for students
    rows = 7
    cols = 3
    table_shape = slide_6.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.5))
    table = table_shape.table
    
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(2.8)
    table.columns[2].width = Inches(6.4)
    
    headers = ["Column Name", "Data Type & Constraints", "Role / Constraints / Validation Context"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(15)
            p.font.color.rgb = GOLD
            p.alignment = PP_ALIGN.CENTER
            
    cols_data = [
        ("id", "INT AUTO_INCREMENT PRIMARY KEY", "Unique numeric identifier key."),
        ("student_id", "VARCHAR(20) NOT NULL UNIQUE", "Unique generated registration code (indexed)."),
        ("email", "VARCHAR(120) NOT NULL UNIQUE", "Unique applicant email (indexed to optimize lookup speed)."),
        ("dob", "DATE NOT NULL", "Candidate date of birth (enforces age validation)."),
        ("major", "VARCHAR(100) NOT NULL", "Selected major program (Computer Science, Data Science, etc.)."),
        ("password_hash", "VARCHAR(255) NOT NULL", "Secure password string encrypted using Bcrypt (rounds: 12).")
    ]
    
    for row_idx, data in enumerate(cols_data, start=1):
        for col_idx, text in enumerate(data):
            cell = table.cell(row_idx, col_idx)
            cell.text = text
            cell.fill.solid()
            if row_idx % 2 == 0:
                cell.fill.fore_color.rgb = LIGHT_GRAY
            else:
                cell.fill.fore_color.rgb = WHITE
                
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = DARK_GRAY
                if col_idx < 2:
                    p.alignment = PP_ALIGN.CENTER
                    p.font.bold = (col_idx == 0)

    # =========================================================================
    # SLIDE 7: Tables & Relationships: System Schema (Tables/Relationships 2/2)
    # =========================================================================
    slide_7 = add_content_slide("Tables & Relationships: System Tables")
    
    # Left table: admins
    tb_left_title = slide_7.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.6), Inches(0.4))
    tb_left_title.text_frame.paragraphs[0].text = "Admins Table Schema"
    tb_left_title.text_frame.paragraphs[0].font.size = Pt(18)
    tb_left_title.text_frame.paragraphs[0].font.bold = True
    tb_left_title.text_frame.paragraphs[0].font.color.rgb = NAVY
    
    table_shape_l = slide_7.shapes.add_table(5, 3, Inches(0.8), Inches(2.1), Inches(5.6), Inches(2.8))
    tl = table_shape_l.table
    tl.columns[0].width = Inches(1.5)
    tl.columns[1].width = Inches(1.8)
    tl.columns[2].width = Inches(2.3)
    
    headers_l = ["Column", "Type", "Details"]
    for i, h in enumerate(headers_l):
        cell = tl.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.text_frame.paragraphs[0].font.color.rgb = GOLD
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(13)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
    admin_cols = [
        ("id", "INT PRIMARY KEY", "PK auto-key"),
        ("username", "VARCHAR(80) UNIQUE", "Unique login username"),
        ("email", "VARCHAR(120) UNIQUE", "Unique administrative contact"),
        ("password_hash", "VARCHAR(255)", "Secure Bcrypt password signature")
    ]
    for row_idx, data in enumerate(admin_cols, start=1):
        for col_idx, text in enumerate(data):
            cell = tl.cell(row_idx, col_idx)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if row_idx % 2 == 0 else WHITE
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY

    # Right table: relationship
    tb_right_title = slide_7.shapes.add_textbox(Inches(6.9), Inches(1.6), Inches(5.6), Inches(0.4))
    tb_right_title.text_frame.paragraphs[0].text = "Dynamic Authorization Schema"
    tb_right_title.text_frame.paragraphs[0].font.size = Pt(18)
    tb_right_title.text_frame.paragraphs[0].font.bold = True
    tb_right_title.text_frame.paragraphs[0].font.color.rgb = NAVY
    
    table_shape_r = slide_7.shapes.add_table(5, 3, Inches(6.9), Inches(2.1), Inches(5.6), Inches(2.8))
    tr = table_shape_r.table
    tr.columns[0].width = Inches(1.6)
    tr.columns[1].width = Inches(1.7)
    tr.columns[2].width = Inches(2.3)
    
    headers_r = ["Concept", "Security Policy", "Session Validation Action"]
    for i, h in enumerate(headers_r):
        cell = tr.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.text_frame.paragraphs[0].font.color.rgb = GOLD
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(13)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
    rel_rows = [
        ("Student session", "HttpOnly Cookie", "Saves regular integer ID in server session."),
        ("Admin session", "HttpOnly Cookie", "Prepends 'admin-' prefix to user loader key."),
        ("Audit Logging", "JSON formatting", "Triggers logs inside audit.py on changes."),
        ("REST Access", "JWT validation", "API routes verify cryptographic tokens.")
    ]
    for row_idx, data in enumerate(rel_rows, start=1):
        for col_idx, text in enumerate(data):
            cell = tr.cell(row_idx, col_idx)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if row_idx % 2 == 0 else WHITE
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY

    # Descriptive text below
    desc_box = slide_7.shapes.add_textbox(Inches(0.8), Inches(5.1), Inches(11.7), Inches(1.5))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    tf_desc.margin_left = tf_desc.margin_top = tf_desc.margin_bottom = tf_desc.margin_right = 0
    p = tf_desc.paragraphs[0]
    p.text = "💡 **Verification & Role Isolation:**"
    p.font.name = "Arial"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = NAVY
    p.space_after = Pt(4)
    
    p = tf_desc.add_paragraph()
    p.text = "• Separation of Concerns: The system checks the Admins table first. If a match is found, session IDs are stored with an 'admin-' prefix, ensuring complete partition from Student accounts.\n• Constraints Enforcement: Database uniqueness constraints prevent duplicated student IDs or candidate emails."
    p.font.name = "Arial"
    p.font.size = Pt(13.5)
    p.font.color.rgb = DARK_GRAY

    # =========================================================================
    # SLIDE 8: SQL Queries: Write Operations (SQL Queries/CRUD Slide 1/2)
    # =========================================================================
    slide_8 = add_content_slide("SQL CRUD: Write Actions (Inserts/Updates)")
    
    left_box = slide_8.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(6.0), Inches(4.8))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    tf_left.margin_left = tf_left.margin_top = tf_left.margin_bottom = tf_left.margin_right = 0
    
    p = tf_left.paragraphs[0]
    p.text = "Student Hub maps write operations dynamically using SQLAlchemy ORM to run secure, SQL-injection-resistant database queries."
    p.font.name = "Arial"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GRAY
    p.space_after = Pt(14)
    
    write_pts = [
        "📥 Candidate Registration INSERT: Runs transaction inserts containing validated personal, address, and graduation details.",
        "🔒 Encrypted Passwords: Computes salt and bcrypt hashes, writing the 255-character string to the database.",
        "✏️ Profile updates: Runs UPDATE statements matching the candidate session ID to alter contact details."
    ]
    for pt in write_pts:
        p = tf_left.add_paragraph()
        p.text = pt
        p.font.name = "Arial"
        p.font.size = Pt(14.5)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(12)

    # Right Box: SQL Code Insert Sample
    right_card = slide_8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.4))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = DARK_GRAY
    right_card.line.color.rgb = GOLD
    right_card.line.width = Pt(1)
    
    rc_text = slide_8.shapes.add_textbox(Inches(7.35), Inches(2.0), Inches(5.0), Inches(4.0))
    tf_rc = rc_text.text_frame
    tf_rc.word_wrap = True
    p = tf_rc.paragraphs[0]
    p.text = "RAW SQL WRITE STATEMENTS"
    p.font.name = "Arial"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.space_after = Pt(12)
    
    sql_write = (
        "-- Registering a student\n"
        "INSERT INTO students (\n"
        "  student_id, first_name, last_name, email, \n"
        "  major, enrollment_type, password_hash\n"
        ") VALUES (\n"
        "  'STU-2026-A3F9', 'John', 'Doe', 'john.doe@email.com',\n"
        "  'Computer Science', 'Full-Time', '$2b$12$R9...'\n"
        ");\n\n"
        "-- Updating candidate profile details\n"
        "UPDATE students SET \n"
        "  phone = '+1 555-0199',\n"
        "  street = '456 Oak Road'\n"
        "WHERE student_id = 'STU-2026-A3F9';"
    )
    p = tf_rc.add_paragraph()
    p.text = sql_write
    p.font.name = "Courier New"
    p.font.size = Pt(9.5)
    p.font.color.rgb = SUCCESS_GREEN
    p.line_spacing = 1.15

    # =========================================================================
    # SLIDE 9: SQL Queries: Read & Analytics (SQL Queries/CRUD Slide 2/2)
    # =========================================================================
    slide_9 = add_content_slide("SQL CRUD: Reads & Aggregation Analytics")
    
    left_box = slide_9.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(6.0), Inches(4.8))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    tf_left.margin_left = tf_left.margin_top = tf_left.margin_bottom = tf_left.margin_right = 0
    
    p = tf_left.paragraphs[0]
    p.text = "The administrative dashboard processes dynamically filtered SELECT queries to support real-time searches and stats summaries."
    p.font.name = "Arial"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GRAY
    p.space_after = Pt(14)
    
    read_pts = [
        "🔍 Roster Search & Sorting: Uses multiple SQL `OR` conditions and dynamic column sorting variables.",
        "📊 Dynamic Aggregations: Groups candidate registries by selected major to render charts.",
        "🔐 Credentials validation: Matches username and loads password hashes to evaluate logins."
    ]
    for pt in read_pts:
        p = tf_left.add_paragraph()
        p.text = pt
        p.font.name = "Arial"
        p.font.size = Pt(14.5)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(12)

    # Right Box: SQL Code Read Sample
    right_card = slide_9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.4))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = DARK_GRAY
    right_card.line.color.rgb = GOLD
    right_card.line.width = Pt(1)
    
    rc_text = slide_9.shapes.add_textbox(Inches(7.35), Inches(2.0), Inches(5.0), Inches(4.0))
    tf_rc = rc_text.text_frame
    tf_rc.word_wrap = True
    p = tf_rc.paragraphs[0]
    p.text = "RAW SQL READ & STATS STATEMENTS"
    p.font.name = "Arial"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.space_after = Pt(12)
    
    sql_read = (
        "-- Search and Sort candidate table\n"
        "SELECT * FROM students \n"
        "WHERE (first_name LIKE '%John%' \n"
        "   OR email LIKE '%John%'\n"
        "   OR student_id = 'John')\n"
        "ORDER BY registration_date DESC;\n\n"
        "-- Aggregations for Chart.js dashboard\n"
        "SELECT major, COUNT(*) AS count \n"
        "FROM students \n"
        "GROUP BY major;\n"
    )
    p = tf_rc.add_paragraph()
    p.text = sql_read
    p.font.name = "Courier New"
    p.font.size = Pt(9.5)
    p.font.color.rgb = SUCCESS_GREEN
    p.line_spacing = 1.15

    # =========================================================================
    # SLIDE 10: Frontend Design System (Frontend Slide 1/2)
    # =========================================================================
    slide_10 = add_content_slide("Frontend UI: Design System & Styling")
    
    left_box = slide_10.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(7.5), Inches(4.8))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    tf_left.margin_left = tf_left.margin_top = tf_left.margin_bottom = tf_left.margin_right = 0
    
    p = tf_left.paragraphs[0]
    p.text = "Student Hub features a premium, responsive frontend design system inspired by the Apple Human Interface Guidelines."
    p.font.name = "Arial"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GRAY
    p.space_after = Pt(14)
    
    design_pts = [
        "🎨 Color Palette: Navy Blue (`#00317e`) representing authority and Gold (`#D4AF37`) representing academic excellence.",
        "🌓 Seamless Dark Mode: Implements dynamic dark styles via `[data-theme='dark']` and saves preference in `localStorage`.",
        "📱 Modern CSS Layouts: Fluid typography clamping, flexbox layouts, glassmorphism cards, and responsive viewports.",
        "♿ Accessibility Touch Targets: Minimum `44px x 44px` targets throughout for easy mobile tap interactions."
    ]
    for pt in design_pts:
        p = tf_left.add_paragraph()
        p.text = pt
        p.font.name = "Arial"
        p.font.size = Pt(14.5)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(10)
        
    # Right Box: CSS Variables
    right_card = slide_10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.4))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = LIGHT_GRAY
    right_card.line.color.rgb = GOLD
    right_card.line.width = Pt(1.5)
    
    rc_text = slide_10.shapes.add_textbox(Inches(9.1), Inches(2.1), Inches(3.1), Inches(3.8))
    tf_rc = rc_text.text_frame
    tf_rc.word_wrap = True
    p = tf_rc.paragraphs[0]
    p.text = "CSS VARIABLE TOKENS"
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.space_after = Pt(14)
    
    tokens = [
        ("--navy", "#00317e (Brand primary)"),
        ("--gold", "#D4AF37 (Brand accent)"),
        ("--card-bg", "rgba(255,255,255,0.65)"),
        ("--card-border", "rgba(255,255,255,0.45)"),
        ("--radius-xl", "24px (Soft pill buttons)")
    ]
    for t_name, t_val in tokens:
        p = tf_rc.add_paragraph()
        p.text = t_name
        p.font.name = "Courier New"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = GOLD
        
        p = tf_rc.add_paragraph()
        p.text = t_val
        p.font.name = "Arial"
        p.font.size = Pt(11.5)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)

    # =========================================================================
    # SLIDE 11: Interface Wizard Screens (Frontend Slide 2/2)
    # =========================================================================
    slide_11 = add_content_slide("Frontend UI: Multi-Step Registration Wizard")
    
    steps = [
        ("Step 1: Personal Info", "First/Last Name, Date of Birth, Gender, Email checks, and formatted Phone inputs.", Inches(0.8)),
        ("Step 2: Mailing Address", "Street address, City, State, ZIP code, and Country drop-down configurations.", Inches(3.8)),
        ("Step 3: Academic Details", "High school name, Graduation year, Major selection, and Enrollment status.", Inches(6.8)),
        ("Step 4: Create Account", "Password inputs with a real-time strength checker and terms acceptance checkboxes.", Inches(9.8))
    ]
    
    for title, desc, left in steps:
        card = slide_11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(2.75), Inches(3.6))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = NAVY
        card.line.width = Pt(1)
        
        tb = slide_11.shapes.add_textbox(left + Inches(0.1), Inches(2.0), Inches(2.55), Inches(3.2))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Arial"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(12)
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.name = "Arial"
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.CENTER

    # Additional text at bottom
    tb_bottom = slide_11.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.2))
    tf_bot = tb_bottom.text_frame
    tf_bot.word_wrap = True
    tf_bot.margin_left = tf_bot.margin_top = tf_bot.margin_bottom = tf_bot.margin_right = 0
    p = tf_bot.paragraphs[0]
    p.text = "💡 **User Experience Logic:**"
    p.font.name = "Arial"
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.space_after = Pt(4)
    
    p = tf_bot.add_paragraph()
    p.text = "• Multi-Step controller: Uses vanilla ES6 JS state controls to hide/show panels without page refreshes.\n• Smart Error Redirects: If registration validation fails, the script identifies error panels and auto-scrolls to the first error field."
    p.font.name = "Arial"
    p.font.size = Pt(13.5)
    p.font.color.rgb = DARK_GRAY

    # =========================================================================
    # SLIDE 12: Innovation: Local Asset Self-Hosting (Innovation Slide 1/2)
    # =========================================================================
    slide_12 = add_content_slide("Innovation: Local Asset Self-Hosting")
    
    left_box = slide_12.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(7.5), Inches(4.8))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    tf_left.margin_left = tf_left.margin_top = tf_left.margin_bottom = tf_left.margin_right = 0
    
    p = tf_left.paragraphs[0]
    p.text = "To eliminate third-party CDN security leaks and support offline deployments, the system self-hosts all design vectors."
    p.font.name = "Arial"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GRAY
    p.space_after = Pt(20)
    
    asset_pts = [
        "🔒 Asset Isolation: Stores the `Material Symbols Outlined` variable font (`.ttf`) directly inside `/static/fonts` directories.",
        "🛡️ CSP Optimization: Satisfies restrictive Content Security Policies, whitelisting `'self'` for font loads instead of public CDN urls.",
        "🔌 Network Autonomy: Guarantees pages render properly and load instantly even when offline.",
        "🧩 Variable Font Axes: Maintains clean CSS access to variable axis values (fill axes, weight, optical sizing) locally."
    ]
    for pt in asset_pts:
        p = tf_left.add_paragraph()
        p.text = pt
        p.font.name = "Arial"
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(12)

    # Right Card: Asset details
    right_card = slide_12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.4))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = LIGHT_GRAY
    right_card.line.color.rgb = GOLD
    right_card.line.width = Pt(1.5)
    
    rc_text = slide_12.shapes.add_textbox(Inches(9.1), Inches(2.1), Inches(3.1), Inches(3.8))
    tf_rc = rc_text.text_frame
    tf_rc.word_wrap = True
    p = tf_rc.paragraphs[0]
    p.text = "LOCAL ASSET SETUP"
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.space_after = Pt(16)
    
    setup_steps = [
        ("Font Binary", "Stored MaterialSymbolsOutlined.ttf locally in the static folder."),
        ("CSS Mapping", "Added @font-face rules mapping font url resources to '/static/fonts/'."),
        ("CSP Lockdown", "Removed Google Fonts external connections, whitelist only 'self' in CSP.")
    ]
    for step_title, step_desc in setup_steps:
        p = tf_rc.add_paragraph()
        p.text = f"✔ {step_title}"
        p.font.name = "Arial"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = GOLD
        p.space_before = Pt(8)
        
        p = tf_rc.add_paragraph()
        p.text = step_desc
        p.font.name = "Arial"
        p.font.size = Pt(11.5)
        p.font.color.rgb = DARK_GRAY

    # =========================================================================
    # SLIDE 13: Innovation: Security & Session Protection (Innovation Slide 2/2)
    # =========================================================================
    slide_13 = add_content_slide("Innovation: Session Protection & Audit Logs")
    
    left_box = slide_13.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(7.5), Inches(4.8))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    tf_left.margin_left = tf_left.margin_top = tf_left.margin_bottom = tf_left.margin_right = 0
    
    p = tf_left.paragraphs[0]
    p.text = "The application implements strict operational auditing and session controllers to secure user states."
    p.font.name = "Arial"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GRAY
    p.space_after = Pt(20)
    
    log_pts = [
        "🛡️ Role Session Separation: Segregates Admin and Student authentication logic, preventing privilege escalation.",
        "📜 JSON Audit Logging (audit.py): Logs operation payloads, user contexts, Client IPs, and target student IDs.",
        "🧹 Input Sanitizer: Clears input strings of SQL strings and script tags before processing to block XSS.",
        "🔑 Token Authentication API: API endpoints use a `@token_required` decorator for JWT authorization."
    ]
    for pt in log_pts:
        p = tf_left.add_paragraph()
        p.text = pt
        p.font.name = "Arial"
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(12)

    # Right Card: Log Format
    right_card = slide_13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.4))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = DARK_GRAY
    right_card.line.color.rgb = GOLD
    right_card.line.width = Pt(1.5)
    
    rc_text = slide_13.shapes.add_textbox(Inches(9.0), Inches(2.0), Inches(3.3), Inches(4.0))
    tf_rc = rc_text.text_frame
    tf_rc.word_wrap = True
    p = tf_rc.paragraphs[0]
    p.text = "JSON AUDIT LOG SAMPLE"
    p.font.name = "Arial"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.space_after = Pt(10)
    
    log_sample = (
        "{\n"
        '  "timestamp": "2026-05-24T10:45:01Z",\n'
        '  "request_id": "8a32-9b2f",\n'
        '  "user_id": 14,\n'
        '  "action": "UPDATE_PROFILE",\n'
        '  "ip": "127.0.0.1",\n'
        '  "details": {\n'
        '    "student_id": "STU-2026-A3",\n'
        '    "status": "active"\n'
        "  }\n"
        "}"
    )
    p = tf_rc.add_paragraph()
    p.text = log_sample
    p.font.name = "Courier New"
    p.font.size = Pt(10.5)
    p.font.color.rgb = RGBColor(0, 255, 255) # Cyan code text
    p.line_spacing = 1.15

    # =========================================================================
    # SLIDE 14: Output Screens: Student Experience (Output Screens Slide 1/2)
    # =========================================================================
    slide_14 = add_content_slide("Output Screens: Student Experience")
    
    outputs_s = [
        ("✨ Hero Particle Canvas", "A custom JavaScript-based floating canvas particle background animates dynamically in the hero viewport, creating a premium look.", Inches(0.8), Inches(1.8)),
        ("⭕ SVG Completion Ring", "The student dashboard displays an interactive SVG completion ring, calculating the percentage of completed profile fields.", Inches(6.8), Inches(1.8)),
        ("🎉 Confetti Celebration", "Triggers a full-screen canvas confetti animation upon successful form submission to welcome new students.", Inches(0.8), Inches(4.3)),
        ("🌓 Dark & Light Modes", "Changes design templates instantly using CSS variables when the dark/light mode toggle is clicked.", Inches(6.8), Inches(4.3))
    ]
    
    for title, desc, left, top in outputs_s:
        card = slide_14.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(5.7), Inches(2.1))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = LIGHT_GRAY
        
        tb = slide_14.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), Inches(5.3), Inches(1.8))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Arial"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.space_after = Pt(6)
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.name = "Arial"
        p.font.size = Pt(12.5)
        p.font.color.rgb = DARK_GRAY

    # =========================================================================
    # SLIDE 15: Output Screens: Admin Console & KPI Charts (Output Screens Slide 2/2)
    # =========================================================================
    slide_15 = add_content_slide("Output Screens: Admin Panel & Analytics")
    
    left_box = slide_15.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(7.5), Inches(4.8))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    tf_left.margin_left = tf_left.margin_top = tf_left.margin_bottom = tf_left.margin_right = 0
    
    p = tf_left.paragraphs[0]
    p.text = "The Admin Console displays key performance metrics and enrollment breakdown visualizations."
    p.font.name = "Arial"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GRAY
    p.space_after = Pt(14)
    
    outputs_a = [
        "📊 KPI Analytics Summary Cards: Displays total student enrollment numbers, active courses, and status counts.",
        "📈 Chart.js Charting: Renders real-time graphs breakdown of academic majors.",
        "👥 Dynamic Searchable Table: Administrators can search entries by ID, email, or name instantly.",
        "📥 CSV Directory Export: Enables downloading the full student database in CSV format with one click."
    ]
    for pt in outputs_a:
        p = tf_left.add_paragraph()
        p.text = pt
        p.font.name = "Arial"
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(10)

    # Right Card: stats JSON representation
    right_card = slide_15.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.4))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = LIGHT_GRAY
    right_card.line.color.rgb = GOLD
    right_card.line.width = Pt(1.5)
    
    rc_text = slide_15.shapes.add_textbox(Inches(9.1), Inches(2.1), Inches(3.1), Inches(3.8))
    tf_rc = rc_text.text_frame
    tf_rc.word_wrap = True
    
    p = tf_rc.paragraphs[0]
    p.text = "STATS API DATA (/api/stats)"
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.space_after = Pt(16)
    
    admin_views = [
        ("JSON Response", "Served dynamically to feed metrics to the landing page counts."),
        ("Majors Breakdown", "Calculates percentage distribution of registered programs."),
        ("CSV Directory Export", "Downloads full candidate registry databases on user request.")
    ]
    for view_title, view_desc in admin_views:
        p = tf_rc.add_paragraph()
        p.text = f"✔ {view_title}"
        p.font.name = "Arial"
        p.font.size = Pt(14.5)
        p.font.bold = True
        p.font.color.rgb = GOLD
        p.space_before = Pt(8)
        
        p = tf_rc.add_paragraph()
        p.text = view_desc
        p.font.name = "Arial"
        p.font.size = Pt(11.5)
        p.font.color.rgb = DARK_GRAY

    # =========================================================================
    # SLIDE 16: Advantages & Limitations
    # =========================================================================
    slide_16 = add_content_slide("Advantages & Limitations")
    
    # Left Card: Advantages
    card_adv = slide_16.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.4))
    card_adv.fill.solid()
    card_adv.fill.fore_color.rgb = LIGHT_GRAY
    card_adv.line.color.rgb = SUCCESS_GREEN
    card_adv.line.width = Pt(1.5)
    
    tb_adv = slide_16.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.2), Inches(4.0))
    tf_adv = tb_adv.text_frame
    tf_adv.word_wrap = True
    p = tf_adv.paragraphs[0]
    p.text = "SYSTEM ADVANTAGES"
    p.font.name = "Arial"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.space_after = Pt(14)
    
    advs = [
        "⭐ 90% Processing Speedup: Reduces active registration times from 30 minutes to under 2.5 minutes.",
        "⭐ Self-Hosted Design System: Serves assets locally, securing page rendering and CSP compliance.",
        "⭐ Parameterized SQL Safety: SQLAlchemy ORM prevents raw injections.",
        "⭐ Clean Administrative Flow: CSV exports and stats dashboards reduce manual staff workloads by ~70%."
    ]
    for pt in advs:
        p = tf_adv.add_paragraph()
        p.text = pt
        p.font.name = "Arial"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(10)

    # Right Card: Limitations
    card_lim = slide_16.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.4))
    card_lim.fill.solid()
    card_lim.fill.fore_color.rgb = LIGHT_GRAY
    card_lim.line.color.rgb = WARNING_RED
    card_lim.line.width = Pt(1.5)
    
    tb_lim = slide_16.shapes.add_textbox(Inches(7.1), Inches(2.1), Inches(5.2), Inches(4.0))
    tf_lim = tb_lim.text_frame
    tf_lim.word_wrap = True
    p = tf_lim.paragraphs[0]
    p.text = "SYSTEM LIMITATIONS"
    p.font.name = "Arial"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.space_after = Pt(14)
    
    lims = [
        "⚠️ Browser State Cookie dependency: Session attributes are lost if cookies are cleared.",
        "⚠️ JavaScript dependent: Client validation checks and UI transitions require JS to be enabled.",
        "⚠️ SQLite single-instance: Default SQLite setup is limited to local developer environments.",
        "⚠️ Multi-Institution Restriction: Lacks support for multi-tenancy configurations."
    ]
    for pt in lims:
        p = tf_lim.add_paragraph()
        p.text = pt
        p.font.name = "Arial"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(10)

    # =========================================================================
    # SLIDE 17: Conclusion
    # =========================================================================
    slide_17 = add_content_slide("Conclusion")
    
    left_box = slide_17.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(7.5), Inches(4.8))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    tf_left.margin_left = tf_left.margin_top = tf_left.margin_bottom = tf_left.margin_right = 0
    
    p = tf_left.paragraphs[0]
    p.text = "Student Hub successfully digitizes university registration processes, providing a secure, responsive, and audit-compliant application."
    p.font.name = "Arial"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GRAY
    p.space_after = Pt(20)
    
    conc_bullets = [
        "🏆 SMART Objectives Achieved: Reached an average completion speed of 2 mins and 31 seconds with zero collisions.",
        "🔒 Enterprise Security Stack: bcrypt password hashing with work factor 12, WTForms CSRF token protection, and rate limiters.",
        "🎨 Optimized UX Layouts: Fluid mobile responsive design compliant with WCAG 2.1 touch guidelines.",
        "⚙️ Deployment Ready: Includes automated run scripts and Docker containerization configuration."
    ]
    for pt in conc_bullets:
        p = tf_left.add_paragraph()
        p.text = pt
        p.font.name = "Arial"
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(12)

    # Right Card: Key Accomplishments
    right_card = slide_17.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.4))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = LIGHT_GRAY
    right_card.line.color.rgb = GOLD
    right_card.line.width = Pt(1.5)
    
    rc_text = slide_17.shapes.add_textbox(Inches(9.1), Inches(2.1), Inches(3.1), Inches(3.8))
    tf_rc = rc_text.text_frame
    tf_rc.word_wrap = True
    p = tf_rc.paragraphs[0]
    p.text = "PROJECT METRICS"
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.space_after = Pt(20)
    
    metrics = [
        ("90% Speedup", "Registration time cut down from 30+ minutes to 2.5 minutes."),
        ("100% Local", "Self-hosted fonts remove third-party tracking risks."),
        ("Compliance logs", "JSON formats support FERPA and GDPR auditing.")
    ]
    for t_title, t_desc in metrics:
        p = tf_rc.add_paragraph()
        p.text = f"✔ {t_title}"
        p.font.name = "Arial"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = GOLD
        p.space_before = Pt(8)
        
        p = tf_rc.add_paragraph()
        p.text = t_desc
        p.font.name = "Arial"
        p.font.size = Pt(11.5)
        p.font.color.rgb = DARK_GRAY

    # =========================================================================
    # SLIDE 18: Future Scope
    # =========================================================================
    slide_18 = add_content_slide("Future Scope & Extensions")
    
    left_box = slide_18.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(7.5), Inches(4.8))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    tf_left.margin_left = tf_left.margin_top = tf_left.margin_bottom = tf_left.margin_right = 0
    
    p = tf_left.paragraphs[0]
    p.text = "Future iterations will build on this secure foundation with integrations and scaling enhancements."
    p.font.name = "Arial"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GRAY
    p.space_after = Pt(20)
    
    future_pts = [
        "💳 Stripe / PayPal Integration: Connects payment gateways to collect registration fees securely.",
        "✉️ SendGrid Email Dispatches: Sends confirmation logs and password reset links to applicants.",
        "🔐 Multi-Factor Authentication (2FA): Adds 2FA options using Google Authenticator codes.",
        " distributed rate limiters: Replaces in-memory limiters with Redis for multi-instance deployments."
    ]
    for pt in future_pts:
        p = tf_left.add_paragraph()
        p.text = pt
        p.font.name = "Arial"
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(12)

    # Right Card: Roadmap
    right_card = slide_18.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.4))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = LIGHT_GRAY
    right_card.line.color.rgb = GOLD
    right_card.line.width = Pt(1.5)
    
    rc_text = slide_18.shapes.add_textbox(Inches(9.1), Inches(2.1), Inches(3.1), Inches(3.8))
    tf_rc = rc_text.text_frame
    tf_rc.word_wrap = True
    p = tf_rc.paragraphs[0]
    p.text = "ROADMAP TIMELINE"
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.space_after = Pt(20)
    
    milestones = [
        ("Short Term", "Email confirmations, password resets, and file uploads."),
        ("Medium Term", "Stripe payment integration and Redis rate limit caching."),
        ("Long Term", "SSO integration, React Native app wrapper, and multi-tenancy.")
    ]
    for m_title, m_desc in milestones:
        p = tf_rc.add_paragraph()
        p.text = f"• {m_title}"
        p.font.name = "Arial"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = GOLD
        p.space_before = Pt(8)
        
        p = tf_rc.add_paragraph()
        p.text = m_desc
        p.font.name = "Arial"
        p.font.size = Pt(11.5)
        p.font.color.rgb = DARK_GRAY

    # =========================================================================
    # SLIDE 19: References
    # =========================================================================
    slide_19 = add_content_slide("References")
    
    left_box = slide_19.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.8))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    tf_left.margin_left = tf_left.margin_top = tf_left.margin_bottom = tf_left.margin_right = 0
    
    p = tf_left.paragraphs[0]
    p.text = "The development and security architecture of Student Hub references the following standards and frameworks:"
    p.font.name = "Arial"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GRAY
    p.space_after = Pt(20)
    
    refs = [
        "📚 **Flask 3.0 Documentation**: https://flask.palletsprojects.com (Framework guides)",
        "🗄️ **SQLAlchemy 2.0 Specifications**: https://www.sqlalchemy.org (Object Relational Mapping)",
        "🎨 **Apple Human Interface Guidelines (HIG)**: https://developer.apple.com/design (UI patterns)",
        "🎨 **Google Material Design 3 Documentation**: https://m3.material.io (UI standards)",
        "🛡️ **OWASP Top 10 Security Risks & Controls**: https://owasp.org (CSRF, injections, XSS checks)",
        "💻 **Werkzeug WSGI Web Server Utilities**: https://werkzeug.palletsprojects.com (Development web server)"
    ]
    for ref in refs:
        p = tf_left.add_paragraph()
        p.text = ref
        p.font.name = "Arial"
        p.font.size = Pt(15.5)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(12)

    # Save presentation
    pptx_path = "Student_Registration_Presentation.pptx"
    prs.save(pptx_path)
    return pptx_path

if __name__ == "__main__":
    path = create_presentation()
    print(f"✅ PowerPoint presentation (19 slides) created successfully: {path}")
    print(f"📍 Location: {os.path.abspath(path)}")
